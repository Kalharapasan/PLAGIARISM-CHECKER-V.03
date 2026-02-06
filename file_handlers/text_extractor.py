import os
import re
import json
import mimetypes
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any, Union
from datetime import datetime
import tempfile
import warnings

class TextExtractor:
    SUPPORTED_FORMATS = {
        '.txt': 'text/plain',
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.doc': 'application/msword',
        '.pdf': 'application/pdf',
        '.rtf': 'application/rtf',
        '.odt': 'application/vnd.oasis.opendocument.text',
        '.html': 'text/html',
        '.htm': 'text/html',
        '.xml': 'application/xml',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.odp': 'application/vnd.oasis.opendocument.presentation',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.ods': 'application/vnd.oasis.opendocument.spreadsheet',
        '.csv': 'text/csv',
        '.py': 'text/x-python',
        '.java': 'text/x-java',
        '.cpp': 'text/x-c++',
        '.c': 'text/x-c',
        '.js': 'application/javascript',
        '.ts': 'application/typescript',
        '.html': 'text/html',
        '.css': 'text/css',
        '.json': 'application/json',
        '.xml': 'application/xml',
        '.md': 'text/markdown',
        '.sql': 'application/sql',
        '.tex': 'application/x-tex',
        '.epub': 'application/epub+zip',
        '.mobi': 'application/x-mobipocket-ebook',
        '.azw': 'application/vnd.amazon.ebook'
       
    }
    
    
    
    def _initialize_handlers(self):
        self.handlers = {}
        self.handlers['.txt'] = self._extract_text_file
        self.handlers['.docx'] = self._extract_docx
        self.handlers['.pdf'] = self._extract_pdf
        self.handlers['.html'] = self._extract_html
        self.handlers['.htm'] = self._extract_html
        self.handlers['.csv'] = self._extract_csv
        self.handlers['.json'] = self._extract_json
        self.handlers['.xml'] = self._extract_xml
        self.handlers['.md'] = self._extract_markdown
        self.handlers['.py'] = self._extract_code
        self.handlers['.java'] = self._extract_code
        self.handlers['.cpp'] = self._extract_code
        self.handlers['.c'] = self._extract_code
        self.handlers['.js'] = self._extract_code
        self.handlers['.ts'] = self._extract_code
        self.handlers['.css'] = self._extract_code
        self.handlers['.sql'] = self._extract_code
        self._register_external_handlers()
    
    def _register_external_handlers(self):
        try:
            from .docx_handler import DOCXHandler
            self.handlers['.docx'] = lambda f: DOCXHandler(self.config).extract_text(f)
            self.handlers['.doc'] = self._extract_doc
        except ImportError:
            pass
        
        try:
            from .pdf_handler import PDFHandler
            self.handlers['.pdf'] = lambda f: PDFHandler(self.config).extract_text(f)
        except ImportError:
            pass
        
        try:
            import striprtf
            self.handlers['.rtf'] = self._extract_rtf
        except ImportError:
            pass
        
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            self.handlers['.odt'] = self._extract_odt
            self.handlers['.odp'] = self._extract_odp
            self.handlers['.ods'] = self._extract_ods
        except ImportError:
            pass
        
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            self.handlers['.epub'] = self._extract_epub
        except ImportError:
            pass
        
        try:
            from pptx import Presentation
            self.handlers['.pptx'] = self._extract_pptx
            self.handlers['.ppt'] = self._extract_ppt
        except ImportError:
            pass
        
        try:
            import pandas as pd
            self.handlers['.xlsx'] = self._extract_xlsx
            self.handlers['.xls'] = self._extract_xls
        except ImportError:
            pass
    
    def _setup_mime_types(self):
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx')
        mimetypes.add_type('application/vnd.ms-powerpoint', '.ppt')
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx')
        mimetypes.add_type('application/vnd.ms-excel', '.xls')
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', '.xlsx')
        mimetypes.add_type('application/vnd.oasis.opendocument.text', '.odt')
        mimetypes.add_type('application/vnd.oasis.opendocument.presentation', '.odp')
        mimetypes.add_type('application/vnd.oasis.opendocument.spreadsheet', '.ods')
        mimetypes.add_type('application/epub+zip', '.epub')
        mimetypes.add_type('application/x-mobipocket-ebook', '.mobi')
        mimetypes.add_type('application/vnd.amazon.ebook', '.azw')
        
    def is_supported_format(self, filepath: str) -> bool:
        ext = Path(filepath).suffix.lower()
        return ext in self.SUPPORTED_FORMATS
        
    def get_file_info(self, filepath: str) -> Dict[str, Any]:
        path = Path(filepath)
        
        if not path.exists():
            raise FileNotFoundError(f"File not found: {filepath}")
        
        mime_type, _ = mimetypes.guess_type(filepath)
        
        info = {
            'filename': path.name,
            'extension': path.suffix.lower(),
            'filepath': str(path.absolute()),
            'file_size': path.stat().st_size,
            'created': datetime.fromtimestamp(path.stat().st_ctime).isoformat(),
            'modified': datetime.fromtimestamp(path.stat().st_mtime).isoformat(),
            'mime_type': mime_type or 'application/octet-stream',
            'is_supported': self.is_supported_format(filepath),
            'is_readable': os.access(filepath, os.R_OK)
        }
        
        return info
    
    def extract_text(self, filepath: str, format_hint: str = None) -> str:
        file_info = self.get_file_info(filepath)
        
        if not file_info['is_supported']:
            if format_hint and format_hint.lower() in self.SUPPORTED_FORMATS:
                extension = format_hint.lower()
            else:
                raise ValueError(f"Unsupported file format: {file_info['extension']}")
        else:
            extension = file_info['extension']
        
        if extension not in self.handlers:
            return self._extract_fallback(filepath)

        try:
            text = self.handlers[extension](filepath)
            text = self._clean_extracted_text(text)
            
            return text
        except Exception as e:
            try:
                text = self._extract_fallback(filepath)
                text = self._clean_extracted_text(text)
                return text
            except Exception as fallback_error:
                raise Exception(f"Failed to extract text from {filepath}: {e}. Fallback also failed: {fallback_error}")
    def extract_with_metadata(self, filepath: str) -> Dict[str, Any]:
        file_info = self.get_file_info(filepath)
        
        try:
            text = self.extract_text(filepath)
            
            result = {
                'success': True,
                'text': text,
                'metadata': file_info,
                'text_metrics': {
                    'characters': len(text),
                    'words': len(text.split()),
                    'lines': text.count('\n') + 1,
                    'paragraphs': len([p for p in text.split('\n\n') if p.strip()])
                },
                'extraction_timestamp': datetime.now().isoformat()
            }    
            format_metadata = self._extract_format_metadata(filepath, file_info['extension'])
            if format_metadata:
                result['format_metadata'] = format_metadata
            
        except Exception as e:
            result = {
                'success': False,
                'error': str(e),
                'metadata': file_info,
                'extraction_timestamp': datetime.now().isoformat()
            }
        
        return result
    
    def extract_from_directory(self, directory: str, 
                              recursive: bool = False,
                              extensions: List[str] = None) -> List[Dict[str, Any]]:
         
        results = []
        dir_path = Path(directory)
        
        if not dir_path.exists() or not dir_path.is_dir():
            raise ValueError(f"Directory not found: {directory}")
        
        if extensions:
            patterns = [f"*{ext}" for ext in extensions]
        else:
            patterns = [f"*{ext}" for ext in self.SUPPORTED_FORMATS.keys()]
        
        files = []
        for pattern in patterns:
            if recursive:
                files.extend(dir_path.rglob(pattern))
            else:
                files.extend(dir_path.glob(pattern))
        files = sorted(set(files))
        
        for file_path in files:
            if file_path.is_file():
                try:
                    result = self.extract_with_metadata(str(file_path))
                    results.append(result)
                except Exception as e:
                    results.append({
                        'success': False,
                        'error': str(e),
                        'filepath': str(file_path),
                        'extraction_timestamp': datetime.now().isoformat()
                    })
        
        return results
    
    
    def _extract_text_file(self, filepath: str) -> str:
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        with open(filepath, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    
    def _extract_docx(self, filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            
            return '\n'.join(text_parts)
        except ImportError:
            return self._extract_docx_fallback(filepath)
    
    def _extract_docx_fallback(self, filepath: str) -> str:
        import zipfile
        import xml.etree.ElementTree as ET
        
        text_parts = []
        with zipfile.ZipFile(filepath) as docx:
            if 'word/document.xml' in docx.namelist():
                xml_content = docx.read('word/document.xml').decode('utf-8')
                text = re.sub(r'<[^>]+>', ' ', xml_content)
                text = re.sub(r'\s+', ' ', text)
                import html
                text = html.unescape(text)
                
                text_parts.append(text.strip())
        
        return '\n'.join(text_parts)
    
    def _extract_doc(self, filepath: str) -> str:
        try:
            import subprocess
            if os.name == 'nt':  
                result = subprocess.run(
                    ['where', 'antiword'],
                    capture_output=True,
                    text=True
                )
                if result.returncode != 0:
                    raise Exception("antiword not found")
                
                antiword_cmd = 'antiword'
            else:  
                antiword_cmd = 'antiword'
            
            result = subprocess.run(
                [antiword_cmd, filepath],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0:
                return result.stdout
            
        except Exception as e:
            print(f"Warning: antiword failed: {e}")
        return self._extract_fallback(filepath)
    
    def _extract_pdf(self, filepath: str) -> str:
        methods = ['pypdf', 'pdfplumber', 'pdfminer', 'fallback']
        
        for method in methods:
            try:
                if method == 'pypdf':
                    from pypdf import PdfReader
                    with open(filepath, 'rb') as f:
                        reader = PdfReader(f)
                        text_parts = []
                        for page in reader.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text_parts.append(page_text)
                        return '\n'.join(text_parts)
                
                elif method == 'pdfplumber':
                    import pdfplumber
                    with pdfplumber.open(filepath) as pdf:
                        text_parts = []
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text_parts.append(page_text)
                        return '\n'.join(text_parts)
                
                elif method == 'pdfminer':
                    from pdfminer.high_level import extract_text
                    return extract_text(filepath)
                
                elif method == 'fallback':
                    import subprocess
                    result = subprocess.run(
                        ['pdftotext', filepath, '-'],
                        capture_output=True,
                        text=True,
                        timeout=30
                    )
                    if result.returncode == 0:
                        return result.stdout
            
            except ImportError:
                continue
            except Exception as e:
                print(f"Warning: PDF extraction method {method} failed: {e}")
                continue
        raise Exception("All PDF extraction methods failed")
    
    def _extract_html(self, filepath: str) -> str:
        try:
            from bs4 import BeautifulSoup
            with open(filepath, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                for script in soup(["script", "style"]):
                    script.decompose()
                    text = soup.get_text()
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    text = '\n'.join(chunk for chunk in chunks if chunk)
                    return text
        except ImportError:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
                text = re.sub(r'<[^>]+>', ' ', content)
                import html
                text = html.unescape(text)
                text = re.sub(r'\s+', ' ', text)
                return text.strip()
    
    def _extract_csv(self, filepath: str) -> str:
        try:
            import pandas as pd
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(filepath, encoding=encoding, on_bad_lines='skip')
                    text_parts = []
                    text_parts.append(', '.join(df.columns.tolist()))
                    for _, row in df.head(1000).iterrows():
                        text_parts.append(', '.join(str(val) for val in row.tolist()))
                    
                    return '\n'.join(text_parts)
                    
                except UnicodeDecodeError:
                    continue
            return self._extract_text_file(filepath)
        except ImportError:
            return self._extract_text_file(filepath)
    
    def _extract_json(self, filepath: str) -> str:
        with open(filepath, 'r', encoding='utf-8') as f:
            try:
                data = json.load(f)
                if isinstance(data, dict):
                    text_parts = self._dict_to_text(data)
                elif isinstance(data, list):
                    text_parts = []
                    for i, item in enumerate(data):
                        text_parts.append(f"Item {i+1}:")
                        if isinstance(item, dict):
                            text_parts.extend(self._dict_to_text(item))
                        else:
                            text_parts.append(str(item))
                else:
                    text_parts = [str(data)]
                
                return '\n'.join(text_parts)
            except json.JSONDecodeError:
                return self._extract_text_file(filepath)
        
    def _extract_xml(self, filepath: str) -> str:
        try:
            import xml.etree.ElementTree as ET
            
            tree = ET.parse(filepath)
            root = tree.getroot()
            text_parts = []
            def extract_element_text(element, depth=0):
                if element.text and element.text.strip():
                    indent = '  ' * depth
                    text_parts.append(f"{indent}{element.tag}: {element.text.strip()}")
                for child in element:
                    extract_element_text(child, depth + 1)
                if element.tail and element.tail.strip():
                    indent = '  ' * depth
                    text_parts.append(f"{indent}Tail: {element.tail.strip()}")
            
            extract_element_text(root)
            return '\n'.join(text_parts)
        except Exception:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
                text = re.sub(r'<[^>]+>', ' ', content)
                import html
                text = html.unescape(text)
                text = re.sub(r'\s+', ' ', text)
                return text.strip()
    
    def _extract_markdown(self, filepath: str) -> str:
        text = self._extract_text_file(filepath)
        text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'\*\*|\*|__|_', '', text)
        text = re.sub(r'```.*?```', '', text, flags=re.DOTALL)
        text = re.sub(r'`.*?`', '', text)
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        text = re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', text)
        text = re.sub(r'^>\s+', '', text, flags=re.MULTILINE)
        text = re.sub(r'^[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)
        
        return text.strip()
    
    def _extract_code(self, filepath: str) -> str:
        text = self._extract_text_file(filepath)
        text = re.sub(r'//.*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'/\*.*?\*/', '', text, flags=re.DOTALL)
        text = re.sub(r'""".*?"""', '', text, flags=re.DOTALL)
        text = re.sub(r"'''.*?'''", '', text, flags=re.DOTALL)
        
        return text.strip()
    
    def _extract_rtf(self, filepath: str) -> str:
        try:
            import striprtf.striprtf as striprtf
            
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
                return striprtf.rtf_to_text(rtf_content)
        except ImportError:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                text = re.sub(r'\\[a-z]+\d*', ' ', content)
                text = re.sub(r'\{[^}]*\}', ' ', text)
                text = re.sub(r'\s+', ' ', text)
                
                return text.strip()
    
    def _extract_odt(self, filepath: str) -> str:
        import zipfile
        import xml.etree.ElementTree as ET
        
        text_parts = []
        with zipfile.ZipFile(filepath) as odt:
            if 'content.xml' in odt.namelist():
                xml_content = odt.read('content.xml').decode('utf-8')
                root = ET.fromstring(xml_content)
                ns = {
                    'office': 'urn:oasis:names:tc:opendocument:xmlns:office:1.0',
                    'text': 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
                }
                for elem in root.findall('.//text:p', ns):
                    if elem.text:
                        text_parts.append(elem.text)
        
        return '\n'.join(text_parts)
    
    def _extract_odp(self, filepath: str) -> str:
        return self._extract_odt(filepath)
    
    def _extract_ods(self, filepath: str) -> str:
        return self._extract_odt(filepath)
    
    def _extract_epub(self, filepath: str) -> str:
        import zipfile
        import xml.etree.ElementTree as ET
        
        text_parts = []
        with zipfile.ZipFile(filepath) as epub:
            if 'META-INF/container.xml' in epub.namelist():
                container_xml = epub.read('META-INF/container.xml').decode('utf-8')
                container_root = ET.fromstring(container_xml)
                ns = {'container': 'urn:oasis:names:tc:opendocument:xmlns:container'}
                rootfile = container_root.find('.//container:rootfile', ns)
                
                if rootfile is not None:
                    opf_path = rootfile.get('full-path')
                    
                    if opf_path in epub.namelist():
                        opf_content = epub.read(opf_path).decode('utf-8')
                        opf_root = ET.fromstring(opf_content)
                        for name in epub.namelist():
                            if name.endswith(('.html', '.xhtml', '.htm')):
                                try:
                                    html_content = epub.read(name).decode('utf-8')
                                    text = re.sub(r'<[^>]+>', ' ', html_content)
                                    text = re.sub(r'\s+', ' ', text)
                                    text_parts.append(text.strip())
                                except:
                                    pass
        
        return '\n'.join(text_parts)
    
    def _extract_pptx(self, filepath: str) -> str:
        try:
            from pptx import Presentation
            
            prs = Presentation(filepath)
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_parts.append(shape.text)
            
            return '\n'.join(text_parts)
        
        except ImportError:
            return self._extract_pptx_fallback(filepath)
        
    def _extract_pptx_fallback(self, filepath: str) -> str:
        import zipfile
        
        text_parts = []
        with zipfile.ZipFile(filepath) as pptx:
            for name in pptx.namelist():
                if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                    try:
                        xml_content = pptx.read(name).decode('utf-8')
                        text = re.sub(r'<[^>]+>', ' ', xml_content)
                        text = re.sub(r'\s+', ' ', text)
                        text_parts.append(text.strip())
                    except:
                        pass
        
        return '\n'.join(text_parts)
    
    def _extract_ppt(self, filepath: str) -> str:
        return self._extract_fallback(filepath)
    
    def _extract_xlsx(self, filepath: str) -> str:
        try:
            import pandas as pd
            xls = pd.ExcelFile(filepath)
            text_parts = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(', '.join(df.columns.tolist()))
                for _, row in df.head(1000).iterrows():
                    text_parts.append(', '.join(str(val) for val in row.tolist()))
                
                text_parts.append('') 
            
            return '\n'.join(text_parts)
        except ImportError:
            return self._extract_xlsx_fallback(filepath)
    
    def _extract_xlsx_fallback(self, filepath: str) -> str:
        import zipfile
        
        text_parts = []
        with zipfile.ZipFile(filepath) as xlsx:
            for name in xlsx.namelist():
                if 'sharedStrings.xml' in name or 'sheet' in name and name.endswith('.xml'):
                    try:
                        xml_content = xlsx.read(name).decode('utf-8')
                        text = re.sub(r'<[^>]+>', ' ', xml_content)
                        text = re.sub(r'\s+', ' ', text)
                        text_parts.append(text.strip())
                    except:
                        pass
        
        return '\n'.join(text_parts)
    
    
    def _extract_xls(self, filepath: str) -> str:
        try:
            import xlrd
            
            workbook = xlrd.open_workbook(filepath)
            text_parts = []
            
            for sheet in workbook.sheets():
                text_parts.append(f"Sheet: {sheet.name}")
                for row in range(min(sheet.nrows, 1000)):
                    row_values = []
                    for col in range(sheet.ncols):
                        cell_value = sheet.cell_value(row, col)
                        row_values.append(str(cell_value))
                    text_parts.append(', '.join(row_values))
                
                text_parts.append('')
            
            return '\n'.join(text_parts)